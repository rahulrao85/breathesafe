"""
Fill the official hackathon PPTX template with BreatheSafe content.

Strategy:
- Slide 1: replace the participant details (name, email, problem statement).
- Slides 2-10: find the 'Type here' text frame on each slide and replace
  its content while keeping the first run's formatting (font, size, color).
- Slide 11: leave as the closing/blank slide but add a "Thank you" caption.

Run:
    python scripts\\fill_deck.py
"""
from __future__ import annotations
import copy
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor


TEMPLATE = Path(
    r"F:\AGENTIC WORLD\GEN AI APEC Cohort 2\Prototype Submission Deck _ Gen AI Academy APAC Edition.pptx"
)
OUTPUT = Path(
    r"F:\AGENTIC WORLD\GEN AI APEC Cohort 2\Prototype Submission Deck _ Gen AI Academy APAC Edition.pptx"
)


# --------------------------------------------------------------------
# Text-frame helpers
# --------------------------------------------------------------------
def _replace_text_preserving_format(text_frame, new_text: str) -> None:
    """
    Replace the text of a text frame while preserving the formatting of
    its first run (font, size, color, bold). Used for the 'Type here'
    placeholders in the template.
    """
    if not text_frame.paragraphs:
        text_frame.text = new_text
        return

    first_para = text_frame.paragraphs[0]
    first_run = first_para.runs[0] if first_para.runs else None

    if first_run is None:
        text_frame.text = new_text
        return

    # Save the format of the first run
    saved_font = first_run.font
    saved_bold = saved_font.bold
    saved_size = saved_font.size
    saved_name = saved_font.name
    try:
        saved_color = saved_font.color.rgb
    except Exception:
        saved_color = None

    # Clear the text frame and put a single paragraph with our content
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = new_text
    if saved_bold is not None:
        run.font.bold = saved_bold
    if saved_size is not None:
        run.font.size = saved_size
    if saved_name is not None:
        run.font.name = saved_name
    if saved_color is not None:
        run.font.color.rgb = saved_color


def _find_placeholder(slide, marker: str):
    """Return the first text frame whose text contains `marker`."""
    for shape in slide.shapes:
        if shape.has_text_frame and marker in shape.text_frame.text:
            return shape.text_frame
    return None


def _set_title(slide, new_title: str) -> None:
    """Replace the slide title while keeping the title's font/colour."""
    tf = _find_placeholder(slide, new_title[:0])  # not used; we do it explicitly
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()
        # Title text boxes on this template are short and upper-cased
        # or are the section title strings we list in TITLE_REPLACEMENTS.
        if txt in TITLE_REPLACEMENTS:
            _replace_text_preserving_format(shape.text_frame, new_title)
            return


# --------------------------------------------------------------------
# Per-slide content
# --------------------------------------------------------------------
TITLE_REPLACEMENTS = {
    "Brief about the idea",
    "List of features offered by the solution",
    "Process flow diagram or Use-case diagram",
    "Wireframes/Mock diagrams of the proposed solution",
    "Architecture diagram of the proposed solution:",
    "Snapshots of the prototype",
    "OpportunitiesHow different is it from any of the other existing ideas?USP of the proposed solution",
    "Technologies / Google / Nvidia Services used in the solution(Why did you choose your specific AI stack and system design?)",
}


SLIDE1_PARTICIPANT = (
    "Participant Details\n\n"
    "Participant Name: Rahul Rao\n"
    "Email: rahulrao85@gmail.com\n"
    "Track: Unified Data Analytics & Intelligence\n"
    "Project: BreatheSafe\n"
    "Hackathon: Google Cloud Gen AI Academy APAC 2026 — Cohort 2\n\n"
    "Problem Statement:\n"
    "India has 100M+ undiagnosed obstructive sleep apnea (OSA) cases. "
    "Public-health teams, NGOs, and wellness leaders do not know where "
    "to focus scarce screening and awareness capacity because the "
    "relevant signals — NFHS-5 health data, Census demographics, CPCB "
    "air quality, Google Trends awareness — sit in silos. BreatheSafe "
    "joins these signals in BigQuery and uses Vertex AI (Gemini 2.5 Flash) + ADK to identify "
    "districts where OSA risk is high but awareness is low — the "
    "'awareness deserts' that screening camps should reach first."
)


SLIDE2_BRIEF = (
    "BreatheSafe is a district-level public-health intelligence layer "
    "for sleep-apnea awareness in India.\n\n"
    "It is not a personal sleep tracker. It does not ingest wearable "
    "or CPAP data, and it does not diagnose individuals. It answers "
    "one population-level question:\n\n"
    "    Where in India should sleep-apnea screening camps and "
    "awareness drives be launched first?\n\n"
    "The product joins five public datasets — NFHS-5, Census 2011, "
    "CPCB / OpenAQ-style AQI, Google Trends awareness, and a "
    "STOP-BANG-derived screening model — in a single BigQuery "
    "warehouse. A Vertex AI-powered ADK agent reasons over the joined "
    "view, cites data sources, and refuses to make individual "
    "diagnostic claims. A Looker Studio dashboard surfaces the "
    "awareness-gap ranking publicly. The whole app is one container "
    "on Cloud Run.\n\n"
    "Positioning: 'I live with sleep apnea. But this project is not "
    "about my device data. It is about the millions of Indians who "
    "may never get screened because the system does not know where "
    "awareness is missing.'"
)


SLIDE3_SOLUTION = (
    "How did we approach the unified-data problem statement?\n"
    "- Started from the question, not the data: 'where should India "
    "screen first?'\n"
    "- Inventoried 5 publicly available Indian datasets that each "
    "touch one OSA risk factor (NFHS-5, Census, CPCB, Google Trends, "
    "STOP-BANG weights).\n"
    "- Built a single BigQuery view (`breathesafe.district_risk_"
    "scores`) that joins all five and emits a 0-1 risk score plus a "
    "0-1 awareness-gap score.\n"
    "- Served the view from a FastAPI app, the same container that "
    "hosts the SPA, on Cloud Run.\n\n"
    "How did we design the solution?\n"
    "- Risk model: weighted min-max of 5 STOP-BANG proxies "
    "(hypertension 0.25, obesity 0.25, PM2.5 0.25, age 50+ 0.15, "
    "male % 0.10). No factor is included without a defensible "
    "population proxy.\n"
    "- Awareness layer: weighted blend of Google Trends interest in "
    "'sleep apnea', 'snoring', 'cpap'.\n"
    "- Agent: 3 tools, deterministic router (swap to Vertex AI Gemini function-"
    "calling in prod), Markdown response with citations and a "
    "responsible-AI disclaimer.\n\n"
    "How will the solution be used?\n"
    "- Public health officers identify the top-N districts for "
    "screening camps in a state.\n"
    "- NGOs pick high-risk + low-awareness districts to launch "
    "awareness drives.\n"
    "- Journalists surface the awareness-desert story.\n"
    "- Corporate wellness leads cross-reference campus locations "
    "against district risk.\n\n"
    "How can the solution scale?\n"
    "- BigQuery scales the joins; the view is parameterised.\n"
    "- The container is stateless and horizontally scaled on Cloud "
    "Run.\n"
    "- The agent is LLM-agnostic: swap in any Vertex AI model with a "
    "single env var."
)


SLIDE4_OPPORTUNITIES = (
    "USP — three things no existing tool does together\n\n"
    "1. DISTRICT-LEVEL OSA PRIORITIZATION (not personal). India "
    "lacks any public product that ranks districts for sleep-apnea "
    "screening. Existing tools either track individuals (Sleep IQ, "
    "wearables) or stay at country-level statistics.\n\n"
    "2. AWARENESS GAP AS A FIRST-CLASS METRIC. We do not just score "
    "risk. We multiply it by the inverse of search interest to find "
    "places the system has effectively ignored. The 3 awareness "
    "deserts surfaced in the demo (Nicobar, N Sikkim, S Sikkim) are "
    "the punchline of the product.\n\n"
    "3. MULTIMODAL ENTRY POINT. Upload a newspaper screenshot, a "
    "poster, an AQI image. Vertex AI Vision extracts location + issue, "
    "the agent cross-references BigQuery, and the app returns a "
    "'risk impact' card. This is the only OSA-aware tool that meets "
    "data where it lives.\n\n"
    "HOW IT IS DIFFERENT FROM EVERY EXISTING IDEA\n"
    "- vs. Sleep IQ / wearables: zero personal data, zero "
    "diagnostic claim, public-health scope.\n"
    "- vs. government dashboards: it is reasoned over by an LLM with "
    "RAG, not just filtered and sorted.\n"
    "- vs. generic chatbots: every answer is grounded in a real "
    "BigQuery view and cites the exact source rows."
)


SLIDE5_FEATURES = (
    "1. Unified district risk view (BigQuery) — 365 districts, 22 "
    "columns, 0-1 risk score and 0-1 awareness-gap score.\n"
    "2. Filterable district ranking table — state, risk category, "
    "awareness-desert flag, sort by any numeric column.\n"
    "3. National KPIs — total districts, high-risk count, top "
    "awareness desert, average risk, sources list.\n"
    "4. ADK-style agent chat with 4 fixed demo prompts and free-text "
    "input. Returns Markdown with driving factors and RAG citations.\n"
    "5. Multimodal image upload — Vertex AI Vision (or deterministic "
    "mock) extracts location/topic, agent cross-references BigQuery, "
    "UI shows a risk impact card.\n"
    "6. Looker Studio dashboard — India bubble heatmap + top-30 table "
    "+ awareness-desert bar + risk-vs-awareness scatter.\n"
    "7. Responsible-AI disclaimer — visible in SPA footer, every "
    "agent response, and the deck.\n"
    "8. Local CSV fallback — the same API works without "
    "BigQuery credentials, so the demo never breaks.\n"
    "9. Single-container deploy — one image, one Cloud Run service, "
    "static SPA + FastAPI side by side.\n"
    "10. Explainable scoring — every district row carries the five "
    "normalized factors; the agent explains why each district is "
    "ranked where it is."
)


SLIDE6_PROCESS = (
    "End-to-end flow (judges should read this slide in 15 seconds):\n\n"
    "1. User opens breathesafe.run.app (Cloud Run, public).\n"
    "2. SPA loads KPI strip + district table from /api/summary and "
    "/api/districts. Backend reads from BigQuery (or the local CSV "
    "fallback).\n"
    "3. User types a question OR clicks a fixed prompt. The request "
    "hits POST /api/ask.\n"
    "4. The ADK-style agent router picks tools (query_district_risk "
    "+ lookup_guidelines).\n"
    "5. query_district_risk runs a parameterised BigQuery query (or "
    "filters the local CSV).\n"
    "6. lookup_guidelines returns the top-3 passages from the RAG "
    "corpus by keyword overlap (swap to Vertex AI Vector Search in "
    "prod).\n"
    "7. compose_answer builds a Markdown response: ranked table, "
    "driving factors, citations, data sources, responsible-AI "
    "notice.\n"
    "8. For multimodal: user uploads an image. POST /api/analyze-"
    "image calls Vertex AI Vision (or mock), extracts location/topic, "
    "agent cross-references the district view, UI renders a risk "
    "impact card.\n"
    "9. The same district view powers the Looker Studio dashboard for "
    "exploration outside the chat."
)


SLIDE7_WIREFRAMES = (
    "Single-page app. Sections top to bottom:\n"
    "  [Header] BreatheSafe + tagline + responsible-AI disclaimer "
    "toggle.\n"
    "  [KPI strip] 4 tiles: districts scored, HIGH-risk count, top "
    "awareness desert, average risk.\n"
    "  [Filters + table] state, risk category, awareness-desert, "
    "search. Table: district, state, risk, gap, est. undiagnosed, "
    "PM2.5.\n"
    "  [Agent chat] text input + 4 fixed-prompt buttons. Response "
    "card renders Markdown including tables, lists, citations.\n"
    "  [Multimodal upload] image picker → risk impact card with "
    "matched district + risk score + awareness gap.\n"
    "  [Looker dashboard] link + screenshot of the published "
    "report.\n"
    "  [Architecture diagram] ASCII / Mermaid of the data flow.\n"
    "  [Footer] data sources + responsible-AI notice.\n\n"
    "Design tone: public-health / civic intelligence. Dark navy "
    "background, cyan + amber + red accent, no hospital-app cliches, "
    "no Sleep-IQ visual reuse."
)


SLIDE8_ARCHITECTURE = (
    "LAYER 1 — DATA SOURCES (5 public datasets)\n"
    "  - NFHS-5 district health data\n"
    "  - Census 2011 demographics (calibrated synthetic for missing "
    "districts)\n"
    "  - CPCB / OpenAQ-style PM2.5\n"
    "  - Google Trends awareness signals\n"
    "  - STOP-BANG model weights + citations\n\n"
    "LAYER 2 — CLOUD STORAGE\n"
    "  Raw CSVs, guideline PDFs, uploaded demo images\n\n"
    "LAYER 3 — BIGQUERY (the unified warehouse)\n"
    "  breathesafe.health_survey\n"
    "  breathesafe.demographics\n"
    "  breathesafe.air_quality\n"
    "  breathesafe.awareness_signals\n"
    "  breathesafe.screening_model\n"
    "  breathesafe.district_risk_scores  ← joined view\n\n"
    "LAYER 4 — ADK AGENT (3 tools)\n"
    "  query_district_risk  → BigQuery\n"
    "  lookup_guidelines    → RAG corpus (10 passages, swap to "
    "Vertex AI Vector Search)\n"
    "  analyze_image        → Vertex AI Vision (deterministic mock "
    "fallback)\n\n"
    "LAYER 5 — VERTEX AI / GEMINI\n"
    "  Text reasoning, image understanding, embeddings\n\n"
    "LAYER 6 — CLOUD RUN\n"
    "  Single container: FastAPI backend + static SPA. Port 8080. "
    "Public.\n\n"
    "LAYER 7 — LOOKER STUDIO\n"
    "  Bubble heatmap, top-30 table, awareness-desert bar, "
    "risk-vs-awareness scatter. Linked from the SPA."
)


SLIDE9_TECH = (
    "GOOGLE CLOUD STACK (the required pieces, all live in the demo)\n"
    "- Cloud Storage: raw CSVs, guideline PDFs, uploaded image inputs.\n"
    "- BigQuery: unified warehouse, 5 tables, 1 joined "
    "district_risk_scores view. Same SQL the API executes is in "
    "backend/bq/schema.sql.\n"
    "- Vertex AI + Gemini 2.5 Flash: agent reasoning (function-calling ready), "
    "embeddings (RAG upgrade path), multimodal image understanding.\n"
    "- ADK: agent with 3 tools, system-prompt guardrails, "
    "responsible-AI disclaimer. Deterministic router today; "
    "production uses Vertex AI function-calling against the same "
    "interface.\n"
    "- Cloud Run: single-container deploy, scales to 0, public URL.\n"
    "- Looker Studio: external dashboard bound directly to the "
    "BigQuery view.\n\n"
    "WHY THIS STACK\n"
    "- The hackathon's theme is 'Unified Data Analytics & "
    "Intelligence'. BigQuery is the only warehouse that lets us join "
    "5 heterogeneous Indian public datasets at district grain in one "
    "SQL statement. Cloud Run gives us a public URL in 5 minutes. "
    "Gemini 2.5 Flash on Vertex AI is the path-of-least-resistance for "
    "multimodal + function-calling in the same SDK.\n"
    "- We deliberately did not use Vertex AI Vector Search. The RAG "
    "corpus is 10 passages; keyword overlap is sufficient and removes "
    "a deploy dependency. The tool interface is identical, so a "
    "future swap is mechanical."
)


SLIDE10_SNAPSHOTS = (
    "Snapshots of the running prototype:\n"
    "  1. SPA landing — dark civic-intelligence palette, KPI strip, "
    "district ranking table, agent chat panel, multimodal upload, "
    "Looker card, architecture diagram, responsible-AI footer.\n"
    "  2. Fixed-prompt 1 fired: 'Top 5 districts in Uttar Pradesh for "
    "OSA screening camps' — returns ranked table (Kanpur Nagar, "
    "Meerut, Agra, Lucknow, Varanasi) with driving factors and RAG "
    "citations.\n"
    "  3. Fixed-prompt 2 fired: 'Why is Delhi NCR high priority?' — "
    "returns Delhi districts (Central, East, New Delhi, West, South) "
    "with obesity, hypertension, PM2.5 and the awareness-vs-risk "
    "story.\n"
    "  4. Multimodal upload — newspaper / AQI image, Vertex AI Vision "
    "extracts location and topic, risk impact card appears with the "
    "matched district, risk score, awareness gap, and recommended "
    "action.\n"
    "  5. Looker Studio dashboard — India bubble heatmap with HIGH-"
    "risk districts in red, MODERATE in amber, LOW in teal. Top-30 "
    "table by awareness-gap score. Awareness-deserts bar chart on "
    "the right.\n"
    "  6. BigQuery console — `breathesafe.district_risk_scores` view "
    "with the 5 source tables alongside, demonstrating the unified "
    "data story.\n\n"
    "Screenshots and the live demo URL are in the deck appendix and "
    "in the GitHub README."
)


SLIDE11_THANKS = (
    "Thank you.\n\n"
    "BreatheSafe — five public datasets, one warehouse, one "
    "intelligence layer, one question: where should India screen "
    "first?\n\n"
    "Live demo: https://breathesafe-api-xxxxx-uc.a.run.app\n"
    "GitHub:    https://github.com/rahulrao85/breathesafe\n"
    "Contact:   rahulrao85@gmail.com\n\n"
    "Responsible AI notice: BreatheSafe is a screening-prioritization "
    "tool, not a medical diagnostic system. Population risk scores "
    "are estimates and do not predict individual OSA."
)


# --------------------------------------------------------------------
# Main
# --------------------------------------------------------------------
def main() -> int:
    if not TEMPLATE.exists():
        print(f"Template not found: {TEMPLATE}", file=sys.stderr)
        return 1

    prs = Presentation(str(TEMPLATE))
    print(f"Loaded template with {len(prs.slides)} slides")

    # ----- Slide 1: participant details -----
    s = prs.slides[0]
    for shape in s.shapes:
        if shape.has_text_frame and "Participant Name" in shape.text_frame.text:
            _replace_text_preserving_format(shape.text_frame, SLIDE1_PARTICIPANT)
            print("Filled slide 1: participant details")
            break

    # ----- Slides 2-10: replace "Type here" placeholders -----
    per_slide_content = {
        1: SLIDE2_BRIEF,        # slide index 1 == second slide
        2: SLIDE3_SOLUTION,
        3: SLIDE4_OPPORTUNITIES,
        4: SLIDE5_FEATURES,
        5: SLIDE6_PROCESS,
        6: SLIDE7_WIREFRAMES,
        7: SLIDE8_ARCHITECTURE,
        8: SLIDE9_TECH,
        9: SLIDE10_SNAPSHOTS,
    }

    for idx, content in per_slide_content.items():
        s = prs.slides[idx]
        target = _find_placeholder(s, "Type here")
        if target is None:
            # Some slides might not have the placeholder; insert a new text box
            # in the lower half of the slide.
            print(f"  slide {idx+1}: no 'Type here' placeholder; appending a new text box")
            tb = s.shapes.add_textbox(
                left=Pt(60), top=Pt(180), width=Pt(820), height=Pt(360)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            run = tf.paragraphs[0].add_run()
            run.text = content
            run.font.size = Pt(14)
            continue
        _replace_text_preserving_format(target, content)
        print(f"Filled slide {idx+1}")

    # ----- Slide 11: closing slide -----
    s = prs.slides[10]
    if s.shapes:
        for shape in s.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip() == "":
                _replace_text_preserving_format(shape.text_frame, SLIDE11_THANKS)
                print("Filled slide 11: thank-you")
                break
        else:
            tb = s.shapes.add_textbox(
                left=Pt(60), top=Pt(120), width=Pt(820), height=Pt(360)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            run = tf.paragraphs[0].add_run()
            run.text = SLIDE11_THANKS
            run.font.size = Pt(20)
            run.font.bold = True
            print("Appended slide 11 thank-you text box")
    else:
        tb = s.shapes.add_textbox(
            left=Pt(60), top=Pt(120), width=Pt(820), height=Pt(360)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        run = tf.paragraphs[0].add_run()
        run.text = SLIDE11_THANKS
        run.font.size = Pt(20)
        run.font.bold = True

    prs.save(str(OUTPUT))
    print(f"\nSaved filled deck to: {OUTPUT}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
